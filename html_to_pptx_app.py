from flask import Flask, request, render_template, send_file
from pptx import Presentation
from bs4 import BeautifulSoup
import os

app = Flask(__name__)

def html_to_pptx(html_content, pptx_file):
    prs = Presentation()
    soup = BeautifulSoup(html_content, 'html.parser')

    for heading in soup.find_all(['h1', 'h2', 'h3']):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = heading.get_text()
        content.text = heading.find_next_sibling(text=True).strip() if heading.find_next_sibling() else ''

    prs.save(pptx_file)

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        file = request.files["html_file"]
        if file:
            html_content = file.read().decode("utf-8")
            output_path = "output.pptx"
            html_to_pptx(html_content, output_path)
            return send_file(output_path, as_attachment=True)
    return render_template("index.html")

if __name__ == "__main__":
    app.run(host='0.0.0.0', port=8080)